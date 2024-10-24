import streamlit as st
import pandas as pd
from sales_by_category import sales_by_category_analysis
from time_slot_analysis import time_slot_analysis
from sales_per_channel import sales_per_channel_analysis
from top_n_brand_sales import top_n_brand_sales_analysis
from brand_availability import top_n_brand_availability_analysis
from top_n_products import top_n_product_analysis
from top_n_product_availability import top_n_product_availability_analysis
from fnb_performance import fnb_performance_analysis
from monetized_brands import analyze_monetized_brands
from counter_shelf_analysis import analyze_counter_shelf_products
from low_performing_brand import low_performing_brand_analysis
from low_performing_products import low_performing_product_analysis
from profit import display_profit_metrics
from pptx import Presentation
from pptx.util import Inches
import numpy as np
import os

# Load your data
data = pd.read_csv("./data/sep_man.csv")
data['orderDate'] = pd.to_datetime(data['orderDate'], format='%d-%m-%Y')

# Set page configuration
st.set_page_config(page_title="TNS DataFactory", page_icon=":bar_chart:", layout="wide")


st.markdown("""
    <style>
    /* Center align the Streamlit container */
    .css-1d391kg {
        display: flex;
        justify-content: center;
    }
    /* Center the page title */
    h1 {
        text-align: center;
        color: #2e7d32; /* Dark green color */
    }
    /* Center the main block */
    section.main {
        max-width: 80%;
        margin-left: auto;
        margin-right: auto;
    }
    /* Box styling for metrics */
    .metrics-box {
        border: 1px solid #ccc;
        border-radius: 10px;
        padding: 10px;
        margin: 10px 0;
        background-color: #f0f4f8; /* Light background color */
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1); /* Subtle shadow */
    }
    /* Strong emphasis text style */
    .metrics-box p {
        font-family: 'Arial', sans-serif; /* Change font */
        font-size: 16px; /* Font size */
        line-height: 1.5; /* Line height */
        margin: 5px 0; /* Margin for paragraphs */
    }
    /* Style for percentage difference */
    .percentage-positive {
        color: green;
        font-weight: bold; /* Make it bold */
    }
    .percentage-negative {
        color: red;
        font-weight: bold; /* Make it bold */
    }

    /* Styles for the Back to Top button */
    .back-to-top {
        position: fixed;
        bottom: 30px;
        right: 30px;
        background-color: #000;
        color: #fff;
        padding: 10px 15px;
        border-radius: 50px;
        text-align: center;
        font-size: 16px;
        cursor: pointer;
        z-index: 9999;
        display: none; /* Initially hidden */
    }
    </style>
    <script>
    // Show button when scrolled down
    window.onscroll = function() {
        var backToTopButton = document.getElementById("backToTop");
        if (document.body.scrollTop > 100 || document.documentElement.scrollTop > 100) {
            backToTopButton.style.display = "block";
        } else {
            backToTopButton.style.display = "none";
        }
    };

    // Scroll to top function
    function scrollToTop() {
        window.scrollTo({
            top: 0,
            behavior: 'smooth' // Smooth scroll
        });
    }
    </script>
""", unsafe_allow_html=True)

# Anchor link at the top of the page
st.markdown("<a id='top'></a>", unsafe_allow_html=True)

# Title Section
st.markdown("<h1>🏭 TNS Data Factory (BETA Mode)</h1>", unsafe_allow_html=True)

# Store selection
store_names = data['storeName'].unique()
selected_store = st.selectbox("Select a Store:", store_names, key="store_selector")

st.markdown(f'<div class="top-bar">Selected Store: {selected_store}</div>', unsafe_allow_html=True)

st.markdown(f"""
    <style>
    .fixed-store-name {{
        position: fixed;
        top: 26px;
        left: 0;
        width: 100%;
        color: green;
        font-size: 28px;
        font-weight: bold;
        padding: 10px;
        z-index: 1000;
        text-align: center;
        border-bottom: 1px solid #ccc;
    }}

    .main {{
        padding-top: 60px;
    }}
    </style>
    <div class="fixed-store-name">{selected_store}</div>
""", unsafe_allow_html=True)


# Date range selection
start_date = st.date_input("Select Start Date:", value=data['orderDate'].min().date(), 
                            min_value=data['orderDate'].min().date(), 
                            max_value=data['orderDate'].max().date())
end_date = st.date_input("Select End Date:", value=data['orderDate'].max().date(), 
                          min_value=data['orderDate'].min().date(), 
                          max_value=data['orderDate'].max().date())

# Convert start_date and end_date to datetime64[ns] for comparison
start_date = pd.to_datetime(start_date)
end_date = pd.to_datetime(end_date)

# Filter data based on selected store and date range
store_data = data[(data['storeName'] == selected_store) & 
                  (data['orderDate'] >= start_date) & 
                  (data['orderDate'] <= end_date)]

# Filter for available stock using the quantity column
store_data_filtered = store_data[store_data['quantity'] > 0]

# Calculate the total number of unique stores for overall data
overall_unique_store_count = data['storeName'].nunique()

# Define the performance rating function based on the logic you provided
def performance_rating(avg_price, overall_average):
    if avg_price > overall_average * 1.2:
        return 'Excellent'
    elif avg_price > overall_average:
        return 'Good'
    elif avg_price == overall_average:
        return 'Average'
    elif avg_price >= overall_average * 0.8:
        return 'Below Average'
    else:
        return 'Poor'

# Calculate total revenue for the selected store
selected_store_total_revenue = store_data['totalProductPrice'].sum()

# Calculate total revenue for all stores
overall_total_revenue = data['totalProductPrice'].sum()

# Calculate average sales for selected store (based on total number of transactions for that store)
selected_store_avg_sales = selected_store_total_revenue / len(store_data)

# Calculate overall average sales using the total number of unique stores (as denominator)
overall_avg_sales = overall_total_revenue / overall_unique_store_count

# Calculate the percentage difference from the overall average
avg_difference_percentage = ((selected_store_avg_sales - overall_avg_sales) 
                             / overall_avg_sales) * 100

# Create a DataFrame for store performance
store_performance = pd.DataFrame({
    'storeName': [selected_store],
    'averageTotalProductPrice': [selected_store_avg_sales],
    'totalRevenue': [selected_store_total_revenue],
    'percentageDifference': [avg_difference_percentage],
    'overallAverage': [overall_avg_sales],
    'overallRevenue': [overall_total_revenue]
})

# Calculate percentage contribution for the selected store
if overall_total_revenue > 0:
    selected_store_percentage_contribution = (selected_store_total_revenue / overall_total_revenue) * 100
else:
    selected_store_percentage_contribution = 0

# Add performance rating column
store_performance['performanceRating'] = store_performance['averageTotalProductPrice'].apply(performance_rating, overall_average=overall_avg_sales)

# Create KPI Cards for key metrics
st.markdown(f"<h2 style='color: green; text-align: center;'>Overall Store KPI: {selected_store}</h2>", unsafe_allow_html=True)

def display_metric(column, label, value, percentage_difference):
    column.metric(label, value)
    # Conditional formatting for delta text
    if percentage_difference > 0:
        delta_text = f'<span style="color: green;">+{percentage_difference:.2f}%</span>'
    elif percentage_difference < 0:
        delta_text = f'<span style="color: red;">{percentage_difference:.2f}%</span>'
    else:
        delta_text = f'<span style="color: grey;">{percentage_difference:.2f}%</span>'
    column.markdown(delta_text, unsafe_allow_html=True)

# Display KPIs as cards
col1, col2, col3 = st.columns(3)

# Store Average Sales
#average_sales = f"₹{store_performance['averageTotalProductPrice'].values[0]:,.2f}"
#percentage_difference = store_performance['percentageDifference'].values[0]
#display_metric(col1, "Store Average Sales", average_sales, percentage_difference)

# Total Revenue
col1.metric("Total Revenue", f"₹{store_performance['totalRevenue'].values[0]:,.2f}")

# Percentage Contribution to Total Revenue
col2.metric("% Contribution to Total Revenue", f"{selected_store_percentage_contribution:.2f}%", delta_color="normal")

# Company Average Sales
col3.metric("Average Sales Overall", f"₹{overall_avg_sales:,.2f}")

# Centering Performance Rating Header
#st.markdown("<h3 style='text-align: center;'>Overall Performance</h3>", unsafe_allow_html=True)

# Get the performance rating
#performance_rating_value = store_performance['performanceRating'].values[0]

# Determine the color based on the performance rating
#if performance_rating_value == 'Excellent':
#    color = 'green'
#elif performance_rating_value in ['Good', 'Average']:
#    color = 'orange'
#else:  # Below Average or Poor
#    color = 'red'

# Centering the Performance Rating Value with color coding
#st.markdown(
#    f"<p style='text-align: center; font-size: 18px; color: {color};'><strong>{performance_rating_value}</strong></p>",
#    unsafe_allow_html=True
#)




selected_store_data = data[data['storeName'] == selected_store]

#hourly
selected_store_data['orderDate'] = pd.to_datetime(selected_store_data['orderDate'], format='%d-%m-%Y')
selected_store_data['time'] = pd.to_datetime(selected_store_data['time'].str.strip(), errors='coerce')

# Extract year and week number from 'orderDate'
selected_store_data['year'] = selected_store_data['orderDate'].dt.isocalendar().year
selected_store_data['week'] = selected_store_data['orderDate'].dt.isocalendar().week


# Extract year and month from 'orderDate'
selected_store_data['year'] = selected_store_data['orderDate'].dt.year
selected_store_data['month'] = selected_store_data['orderDate'].dt.month

# Filter data for the selected date range
filtered_data = data[(data['orderDate'] >= start_date) & (data['orderDate'] <= end_date)]

# Extract month from orderDate
filtered_data['month'] = filtered_data['orderDate'].dt.month

# Calculate monthly sales for each store
monthly_sales_per_store = filtered_data.groupby(['storeName', 'month'])['totalProductPrice'].sum().reset_index()

# Calculate average monthly sales for each store
average_monthly_sales_per_store = monthly_sales_per_store.groupby('storeName')['totalProductPrice'].mean().reset_index()
average_monthly_sales_per_store.rename(columns={'totalProductPrice': 'averageMonthlySales'}, inplace=True)

# Calculate overall average monthly sales based on unique store names in the filtered data
overall_average_monthly_sales = filtered_data['totalProductPrice'].sum() / filtered_data['month'].nunique() / filtered_data['storeName'].nunique()

# Calculate percentage difference
average_monthly_sales_per_store['percentageDifference'] = (
    (average_monthly_sales_per_store['averageMonthlySales'] - overall_average_monthly_sales) / overall_average_monthly_sales
) * 100

# Add overall average to the DataFrame
average_monthly_sales_per_store['overallAverageMonthlySales'] = overall_average_monthly_sales


data['orderDate'] = pd.to_datetime(data['orderDate'], format='%d-%m-%Y', errors='coerce')

# Filter data for the selected date range
filtered_data = data[(data['orderDate'] >= start_date) & (data['orderDate'] <= end_date)]

# Calculate daily sales for each store
daily_sales_per_store = filtered_data.groupby(['storeName', 'orderDate'])['totalProductPrice'].sum().reset_index()

# Calculate average daily sales for each store
average_daily_sales_per_store = daily_sales_per_store.groupby('storeName')['totalProductPrice'].mean().reset_index()
average_daily_sales_per_store.rename(columns={'totalProductPrice': 'averageDailySales'}, inplace=True)

# Calculate overall average daily sales based on unique store names in the filtered data
overall_average_daily_sales = filtered_data['totalProductPrice'].sum() / filtered_data['orderDate'].nunique() / filtered_data['storeName'].nunique()

# Calculate percentage difference
average_daily_sales_per_store['percentageDifference'] = (
    (average_daily_sales_per_store['averageDailySales'] - overall_average_daily_sales) / overall_average_daily_sales
) * 100

# Add overall average daily sales to the DataFrame
average_daily_sales_per_store['overallAverageDailySales'] = overall_average_daily_sales



# hourly calculation
data['orderDate'] = pd.to_datetime(data['orderDate'], format='%d-%m-%Y')  # Convert orderDate to datetime
data['time'] = pd.to_datetime(data['time'], format='%H:%M:%S').dt.time
data['datetime'] = data.apply(lambda row: row['orderDate'] + pd.Timedelta(hours=row['time'].hour,
                                                                        minutes=row['time'].minute,
                                                                        seconds=row['time'].second), axis=1)

# Step 2: Extract hour
data['hour'] = data['datetime'].dt.hour

# Step 3: Filter data for the selected date range
filtered_data = data[(data['orderDate'] >= start_date) & (data['orderDate'] <= end_date)]

# Step 4: Group by store, date, and hour, calculate total sales
hourly_sales_per_store = filtered_data.groupby(['storeName', 'orderDate', 'hour'])['totalProductPrice'].sum().reset_index()

# Step 5: Calculate hourly average sales for each store
hourly_average_sales_per_store = hourly_sales_per_store.groupby(['storeName', 'hour'])['totalProductPrice'].mean().reset_index()
hourly_average_sales_per_store.rename(columns={'totalProductPrice': 'averageHourlySales'}, inplace=True)

# Step 6: Calculate overall average hourly sales based on unique store names for the selected date range
overall_average_hourly_sales = filtered_data['totalProductPrice'].sum() / (filtered_data['storeName'].nunique() * 24)  # Assuming 24 hours in a day

# Step 7: Calculate percentage difference
hourly_average_sales_per_store['percentageDifference'] = (
    (hourly_average_sales_per_store['averageHourlySales'] - overall_average_hourly_sales) / overall_average_hourly_sales
) * 100

# Step 8: Add overall average to the DataFrame
hourly_average_sales_per_store['overallAverageHourlySales'] = overall_average_hourly_sales



# Convert 'orderDate' to datetime
data['orderDate'] = pd.to_datetime(data['orderDate'], format='%d-%m-%Y', errors='coerce')

# Filter data for the selected date range
filtered_data = data[(data['orderDate'] >= start_date) & (data['orderDate'] <= end_date)]

# Extract week number from orderDate
filtered_data['week_number'] = filtered_data['orderDate'].dt.isocalendar().week

# Calculate weekly sales for each store
weekly_sales_per_store = filtered_data.groupby(['storeName', 'week_number'])['totalProductPrice'].sum().reset_index()

# Calculate average weekly sales for each store
average_weekly_sales_per_store = weekly_sales_per_store.groupby('storeName')['totalProductPrice'].mean().reset_index()
average_weekly_sales_per_store.rename(columns={'totalProductPrice': 'averageWeeklySales'}, inplace=True)

# Calculate overall average weekly sales based on unique store names in the filtered data
overall_average_weekly_sales = filtered_data['totalProductPrice'].sum() / filtered_data['week_number'].nunique() / filtered_data['storeName'].nunique()

# Calculate percentage difference
average_weekly_sales_per_store['percentageDifference'] = (
    (average_weekly_sales_per_store['averageWeeklySales'] - overall_average_weekly_sales) / overall_average_weekly_sales
) * 100

# Add overall average to the DataFrame
average_weekly_sales_per_store['overallAverageWeeklySales'] = overall_average_weekly_sales



# Streamlit UI layout
col1, col2, col3, col4 = st.columns([1, 1, 1, 1])

with col1:
    # Stylish boundary box for daily sales KPI
    st.markdown(
        f"""
        <div style='background-color: #f0f0f0; border: 2px solid #0072B8; border-radius: 10px; padding: 8px; width: 250px; margin: auto;'> 

        <h1 style='text-align: center; margin: 0; font-size: 26px;'>Daily:</h1>
        <hr style='border: 1px solid #0072B8; width: 80%; margin: auto;'/>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Store Avg. Daily Sales</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {average_daily_sales_per_store[average_daily_sales_per_store['storeName'] == selected_store]['averageDailySales'].mean():,.2f}</h3>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Overall Avg.</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {overall_average_daily_sales:,.2f}</h3>
        """,
        unsafe_allow_html=True
    )

    # Conditional formatting for percentage difference in daily sales
    percentage_difference_daily = average_daily_sales_per_store[average_daily_sales_per_store['storeName'] == selected_store]['percentageDifference'].mean()
    
    if percentage_difference_daily > 0:
        st.markdown(f"<h4 style='color: green; text-align: center; margin: 0; font-size: 16px;'>+{percentage_difference_daily:.2f}%</h4>", unsafe_allow_html=True)
    elif percentage_difference_daily < 0:
        st.markdown(f"<h4 style='color: red; text-align: center; margin: 0; font-size: 16px;'>{percentage_difference_daily:.2f}%</h4>", unsafe_allow_html=True)
    else:
        st.markdown(f"<h4 style='color: orange; text-align: center; margin: 0; font-size: 16px;'>No Change</h4>", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

# Display a message if no data is available for the selected store and date range
if filtered_data.empty:
    st.markdown("<h4 style='text-align: center; color: red;'>No data available for the selected store and date range.</h4>", unsafe_allow_html=True)

# Streamlit KPI Display Code
with col2:
    # Stylish boundary box for hourly sales KPI
    st.markdown(
        f"""
        <div style='background-color: #f0f0f0; border: 2px solid #0072B8; border-radius: 10px; padding: 8px; width: 250px; margin: auto;'> 

        <h1 style='text-align: center; margin: 0; font-size: 26px;'>Hourly:</h1>
        <hr style='border: 1px solid #0072B8; width: 80%; margin: auto;'/>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Store Avg. Hourly Sales</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {hourly_average_sales_per_store[hourly_average_sales_per_store['storeName'] == selected_store]['averageHourlySales'].mean():,.2f}</h3>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Overall Avg.</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {overall_average_hourly_sales:,.2f}</h3>
        """,
        unsafe_allow_html=True
    )

    # Conditional formatting for percentage difference in hourly sales
    percentage_difference_hourly = hourly_average_sales_per_store[hourly_average_sales_per_store['storeName'] == selected_store]['percentageDifference'].mean()
    
    if percentage_difference_hourly > 0:
        st.markdown(f"<h4 style='color: green; text-align: center; margin: 0; font-size: 16px;'>+{percentage_difference_hourly:.2f}%</h4>", unsafe_allow_html=True)
    elif percentage_difference_hourly < 0:
        st.markdown(f"<h4 style='color: red; text-align: center; margin: 0; font-size: 16px;'>{percentage_difference_hourly:.2f}%</h4>", unsafe_allow_html=True)
    else:
        st.markdown(f"<h4 style='color: orange; text-align: center; margin: 0; font-size: 16px;'>No Change</h4>", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

# Display a message if no data is available for the selected store and date range
if filtered_data.empty:
    st.markdown("<h4 style='text-align: center; color: red;'>No data available for the selected store and date range.</h4>", unsafe_allow_html=True)


# Streamlit KPI Display Code for Weekly Sales
with col3:
    # Stylish boundary box for weekly sales KPI
    st.markdown(
        f"""
        <div style='background-color: #f0f0f0; border: 2px solid #0072B8; border-radius: 10px; padding: 8px; width: 250px; margin: auto;'> 

        <h1 style='text-align: center; margin: 0; font-size: 26px;'>Weekly:</h1>
        <hr style='border: 1px solid #0072B8; width: 80%; margin: auto;'/>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Store Avg. Weekly Sales</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {average_weekly_sales_per_store[average_weekly_sales_per_store['storeName'] == selected_store]['averageWeeklySales'].mean():,.2f}</h3>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Overall Avg.</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {overall_average_weekly_sales:,.2f}</h3>
        """,
        unsafe_allow_html=True
    )

    # Conditional formatting for percentage difference in weekly sales
    percentage_difference_weekly = average_weekly_sales_per_store[average_weekly_sales_per_store['storeName'] == selected_store]['percentageDifference'].mean()
    
    if percentage_difference_weekly > 0:
        st.markdown(f"<h4 style='color: green; text-align: center; margin: 0; font-size: 16px;'>+{percentage_difference_weekly:.2f}%</h4>", unsafe_allow_html=True)
    elif percentage_difference_weekly < 0:
        st.markdown(f"<h4 style='color: red; text-align: center; margin: 0; font-size: 16px;'>{percentage_difference_weekly:.2f}%</h4>", unsafe_allow_html=True)
    else:
        st.markdown(f"<h4 style='color: orange; text-align: center; margin: 0; font-size: 16px;'>No Change</h4>", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

# Display a message if no data is available for the selected store and date range
if filtered_data.empty:
    st.markdown("<h4 style='text-align: center; color: red;'>No data available for the selected store and date range.</h4>", unsafe_allow_html=True)

with col4:  # Monthly KPI
    # Stylish boundary box for monthly sales KPI
    st.markdown(
        f"""
        <div style='background-color: #f0f0f0; border: 2px solid #0072B8; border-radius: 10px; padding: 8px; width: 250px; margin: auto;'> 

        <h1 style='text-align: center; margin: 0; font-size: 26px;'>Monthly:</h1>
        <hr style='border: 1px solid #0072B8; width: 80%; margin: auto;'/>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Store Avg. Monthly Sales</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {average_monthly_sales_per_store[average_monthly_sales_per_store['storeName'] == selected_store]['averageMonthlySales'].mean():,.2f}</h3>

        <h4 style='text-align: center; margin: 0; font-size: 14px;'>Overall Avg.</h4>
        <h3 style='text-align: center; margin: 5px 0; font-size: 24px;'>₹ {overall_average_monthly_sales:,.2f}</h3>
        """,
        unsafe_allow_html=True
    )

    # Conditional formatting for percentage difference in monthly sales
    percentage_difference_monthly = average_monthly_sales_per_store[average_monthly_sales_per_store['storeName'] == selected_store]['percentageDifference'].mean()
    
    if percentage_difference_monthly > 0:
        st.markdown(f"<h4 style='color: green; text-align: center; margin: 0; font-size: 16px;'>+{percentage_difference_monthly:.2f}%</h4>", unsafe_allow_html=True)
    elif percentage_difference_monthly < 0:
        st.markdown(f"<h4 style='color: red; text-align: center; margin: 0; font-size: 16px;'>{percentage_difference_monthly:.2f}%</h4>", unsafe_allow_html=True)
    else:
        st.markdown(f"<h4 style='color: orange; text-align: center; margin: 0; font-size: 16px;'>No Change</h4>", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

# Display a message if no data is available for the selected store and date range
if filtered_data.empty:
    st.markdown("<h4 style='text-align: center; color: red;'>No data available for the selected store and date range.</h4>", unsafe_allow_html=True)



display_profit_metrics(data, selected_store)


# Run all analyses and capture the relevant DataFrames for the report
all_data = data.copy()
sales_by_category_analysis(store_data, all_data)


# Center-align the title, make it green, and add a reasonable icon
st.markdown("---")
st.markdown("<h1 style='color: green; text-align: center;'>⏰ TIME SLOT ANALYSIS</h1>", unsafe_allow_html=True)
# Add separator line below the title
st.markdown("<hr style='border-top: 3px solid #bbb;'>", unsafe_allow_html=True)


st.markdown(
    f"<h3 style='color: green; text-align: center;'>Store KPI from : {start_date.date()} to {end_date.date()}</h3>",
    unsafe_allow_html=True
)

# Function to calculate new metrics
def calculate_additional_metrics(store_data_filtered):
    average_sales_selected = store_data_filtered['totalProductPrice'].sum()/store_data_filtered["storeName"].nunique()
    total_revenue_selected = store_data_filtered['totalProductPrice'].sum()
    
    overall_average_dynamic = data['totalProductPrice'].sum()/data["storeName"].nunique()
    overall_total_revenue_dynamic = data['totalProductPrice'].sum()
    
    selected_store_percentage_contribution_dynamic = (total_revenue_selected / overall_total_revenue_dynamic) * 100 if overall_total_revenue_dynamic > 0 else 0
    percentage_difference_dynamic = ((average_sales_selected - overall_average_dynamic) / overall_average_dynamic) * 100 if overall_average_dynamic > 0 else 0

    return average_sales_selected, total_revenue_selected, selected_store_percentage_contribution_dynamic, overall_average_dynamic, percentage_difference_dynamic

# Get additional metrics
average_sales_selected, total_revenue_selected, selected_store_percentage_contribution_dynamic, overall_average_dynamic, percentage_difference_dynamic = calculate_additional_metrics(store_data_filtered)

# Display additional metrics function
def display_additional_metric(column, label, value, percentage_difference):
    column.metric(label, value)
    if percentage_difference > 0:
        delta_text = f'<span style="color: green;">+{percentage_difference:.2f}%</span>'
    elif percentage_difference < 0:
        delta_text = f'<span style="color: red;">{percentage_difference:.2f}%</span>'
    else:
        delta_text = f'<span style="color: grey;">{percentage_difference:.2f}%</span>'
    column.markdown(delta_text, unsafe_allow_html=True)

# Display additional KPIs as cards
col5, col6, col7, col8, col9 = st.columns(5) 

# Store Average Sales (Dynamic)
display_additional_metric(col5, "Avg Sales", f"₹{total_revenue_selected:,.2f}", percentage_difference_dynamic)

# Percentage Contribution to Total Revenue (Dynamic)
col6.metric("% Contribution to Total Revenue", f"{selected_store_percentage_contribution_dynamic:.2f}%", delta_color="normal")

# Company Average Sales (Dynamic)
col7.metric("Average Sales", f"₹{overall_average_dynamic:,.2f}")

# Difference Between Company and Selected Store Average Sales (Dynamic)
col8.metric("Difference", f"{percentage_difference_dynamic:.2f}%", delta_color="normal")


time_slot_analysis(store_data)


sales_per_channel_analysis(store_data, data)

top_n_brand_df = top_n_brand_sales_analysis(store_data)
top_n_brand_availability_analysis(store_data_filtered)
top_n_product_analysis(store_data)
top_n_product_availability_analysis(store_data_filtered)
fnb_performance_analysis(store_data)
analyze_monetized_brands(store_data)
analyze_counter_shelf_products(store_data)
low_performing_brand_analysis(store_data)
low_performing_product_analysis(store_data)

# Function to create a PowerPoint report
def create_ppt_report():
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TNS Data Factory Report"
    subtitle.text = f"Store: {selected_store}\nDate Range: {start_date.date()} - {end_date.date()}"

    plot_files = [
        'top_n_brands_availability_bar.png',
        'top_n_brands_availability_donut.png',
        'top_n_brands_availability_line.png',
        # more plot WIP
    ]

    for plot_file in plot_files:
        if os.path.exists(plot_file):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = plot_file.replace(".png", "").replace("_", " ").title()
            slide.shapes.add_picture(plot_file, Inches(1), Inches(1), width=Inches(8), height=Inches(4))

    pptx_file = "TNS_Data_Factory_Report.pptx"
    prs.save(pptx_file)
    return pptx_file

# Button to create PowerPoint report
if st.button("Generate PowerPoint Report"):
    pptx_file = create_ppt_report()
    st.success(f"Report created: {pptx_file}")

        # Reading the generated PPTX file in binary mode
    with open(pptx_file, "rb") as file:
        btn = st.download_button(
            label="Download PowerPoint Report",
            data=file,
            file_name=pptx_file,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

st.markdown("<div id='backToTop' class='back-to-top' onclick='scrollToTop()'>Back to Top</div>", unsafe_allow_html=True)

st.markdown("<footer><p style='text-align:center;'>TNS Data Factory © 2024</p></footer>", unsafe_allow_html=True)
